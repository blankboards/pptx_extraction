from pptx import Presentation  
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageEnhance, ImageFilter
import pytesseract
import cv2
import numpy as np
import io
import os

pytesseract.pytesseract.tesseract_cmd = r'C:\\Program Files\\Tesseract-OCR\\tesseract.exe'

def extract_images_from_ppt_tesseract(file_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    image_texts = []
    presentation = Presentation(file_path)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_folder = os.path.join(output_dir, f"slide_{slide_number}", "image")
        os.makedirs(slide_folder, exist_ok=True)
        image_index = 1  

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_name = f"image_{image_index}.jpg"
                image_path = os.path.join(slide_folder, image_name)

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                text = process_image_for_ocr(image_bytes)
                if text:  
                    image_texts.append(f"Slide {slide_number}, Image {image_index} Text:\n{text}")
                image_index += 1  
    return image_texts

def process_image_for_ocr(image_bytes):
    try:
        if not contains_text(image_bytes):
          return ""  # 直接跳过 OCR，避免乱码
        # 读取图片
        with Image.open(io.BytesIO(image_bytes)) as img:
            img = img.convert('L')  # 转换为灰度
            img = ImageEnhance.Contrast(img).enhance(2.0) # 增强对比度
            img_cv = np.array(img)  # 将 PIL 图片转换为 OpenCV 格式
            img_cv = cv2.medianBlur(img_cv, 3)  # 应用中值滤波（减少噪声）
            _, img_cv = cv2.threshold(img_cv, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)  # Otsu 二值化（自动调整阈值）
            edges = cv2.Canny(img_cv, 30, 150) # 边缘检测，检查是否含有文字
            edge_ratio = np.sum(edges > 0) / edges.size
            if edge_ratio < 0.005:  # 设置一个经验阈值（如果边缘像素太少，说明可能没有文字）
                return ""  # 直接跳过 OCR 处理
            text = pytesseract.image_to_string(img_cv, lang='chi_sim+eng', config='--psm 6')  # OCR 识别（使用合适的 PSM 模式）
            # 过滤无效文本（防止全是符号、噪声）
            if len(text.strip()) < 5:  # 设定最小长度阈值
                return ""
            return text.strip()
    except Exception as e:
        return f"OCR processing failed: {e}"

def contains_text(image_bytes):
    """判断图片是否含有文本"""
    with Image.open(io.BytesIO(image_bytes)) as img:
        img_cv = np.array(img.convert('L'))  # 转换为灰度
        img_blur = cv2.GaussianBlur(img_cv, (5, 5), 0) # 应用高斯模糊，减少噪声
        _, img_bin = cv2.threshold(img_blur, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU) # 使用 Otsu 二值化
        edges = cv2.Canny(img_bin, 30, 150)  # 计算边缘像素比例
        edge_ratio = np.sum(edges > 0) / edges.size
        contours, _ = cv2.findContours(img_bin, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)  # 轮廓检测
        text_area = sum(cv2.contourArea(cnt) for cnt in contours)
        text_ratio = text_area / (img_cv.shape[0] * img_cv.shape[1])
        # **文本判断标准**：
        # - **边缘像素比例 > 0.005** (说明有文字结构)
        # - **文本区域面积比 > 0.01** (大于1%图像面积)
        if edge_ratio > 0.005 and text_ratio > 0.01:
            print('word')
            return True  # 说明图像含有文字
    return False  # 图像无明显文字