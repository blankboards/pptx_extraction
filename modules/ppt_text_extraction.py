# ppt_text_extraction.py
from pptx import Presentation 
from modules.config import OUTPUT_DIR_2
import os
import logging
import win32com.client
import pythoncom

# 配置日志（与 app.py 和 main.py 一致）
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

def extract_text_from_ppt(file_path):
    os.makedirs(OUTPUT_DIR_2, exist_ok=True)
    # 判断文件格式
    ext = os.path.splitext(file_path.lower())[1]
    if ext != '.pptx':
        return extract_text_from_ppt_legacy(file_path)
    
    text_output = []
    presentation = Presentation(file_path) 
    logger.info(f"Processing PPTX file: {file_path}")
    
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text.append(run.text)
        if slide_text:
            text_output.append(f"\n\n@@@Slide_{slide_number}@@@\n" + "\n".join(slide_text))
    logger.info(f"Completed processing {file_path}. Extracted text from {len(text_output)} slides.")
    return text_output

def extract_metadata(file_path):
    try:
        presentation = Presentation(file_path)
        props = presentation.core_properties
        logger.info(f"Extracting metadata from PPTX: {file_path}")

        metadata = {
            "Title": props.title or "N/A",
            "Author": props.author or "N/A",
            "Subject": props.subject or "N/A",
            "Keywords": props.keywords or "N/A",
            "Comments": props.comments or "N/A",
            "Last Modified By": props.last_modified_by or "N/A",
            "Created": props.created.strftime('%Y-%m-%d %H:%M:%S') if props.created else "N/A",
            "Modified": props.modified.strftime('%Y-%m-%d %H:%M:%S') if props.modified else "N/A",
            "Category": props.category or "N/A",
            "Content Status": props.content_status or "N/A",
            "Identifier": props.identifier or "N/A",
            "Language": props.language or "N/A",
            "Revision": props.revision or "N/A"
        }
        logger.info(f"Metadata extracted successfully from {file_path}")
        return metadata

    except Exception as e:
        logger.error(f"Error extracting metadata from PPTX {file_path}: {e}")
        return {"Error": str(e)}

def extract_text_from_ppt_legacy(file_path, output_dir=OUTPUT_DIR_2):
    """
    从非 PPTX 格式的文件（如 .ppt, .pot, .pps）中提取文本。

    Args:
        file_path (str): PPT 文件路径。
        output_dir (str): 输出目录，默认为 OUTPUT_DIR_2。

    Returns:
        list: 包含每页幻灯片文本的列表，与 extract_text_from_ppt 输出格式一致。
    """
    try:
        pythoncom.CoInitialize()  # 初始化 COM 环境
        os.makedirs(output_dir, exist_ok=True)
        text_output = []
        app = win32com.client.Dispatch("PowerPoint.Application")
        prs = app.Presentations.Open(file_path, WithWindow=False)
        logger.info(f"Processing legacy PPT format: {file_path}")

        for slide_number, slide in enumerate(prs.Slides, start=1):
            slide_text = []
            for shape in slide.Shapes:
                if shape.HasTextFrame and shape.TextFrame.HasText:
                    text = shape.TextFrame.TextRange.Text.strip()
                    if text:
                        slide_text.append(text)
            if slide_text:
                text_output.append(f"\n\n@@@Slide_{slide_number}@@@\n" + "\n".join(slide_text))

        prs.Close()
        app.Quit()
        pythoncom.CoUninitialize()  # 释放 COM 环境
        logger.info(f"Completed processing {file_path}. Extracted text from {len(text_output)} slides.")
        return text_output

    except Exception as e:
        logger.error(f"Error processing legacy PPT format {file_path}: {e}")
        if 'app' in locals():
            app.Quit()
        pythoncom.CoUninitialize()  # 确保异常时也释放 COM
        return []

def extract_metadata_from_ppt_legacy(file_path):
    """
    从非 PPTX 格式的文件（如 .ppt, .pot, .pps）中提取元数据。

    Args:
        file_path (str): PPT 文件路径。

    Returns:
        dict: 元数据字典，与 extract_metadata 输出格式一致。
    """
    try:
        pythoncom.CoInitialize()  # 初始化 COM 环境
        app = win32com.client.Dispatch("PowerPoint.Application")
        prs = app.Presentations.Open(file_path, WithWindow=False)
        props = prs.BuiltInDocumentProperties
        logger.info(f"Extracting metadata from legacy PPT format: {file_path}")

        metadata = {
            "Title": props("Title").Value if props("Title") else "N/A",
            "Author": props("Author").Value if props("Author") else "N/A",
            "Subject": props("Subject").Value if props("Subject") else "N/A",
            "Keywords": props("Keywords").Value if props("Keywords") else "N/A",
            "Comments": props("Comments").Value if props("Comments") else "N/A",
            "Last Modified By": props("Last Saved By").Value if props("Last Saved By") else "N/A",
            "Created": props("Creation Date").Value.strftime('%Y-%m-%d %H:%M:%S') if props("Creation Date") else "N/A",
            "Modified": props("Last Save Time").Value.strftime('%Y-%m-%d %H:%M:%S') if props("Last Save Time") else "N/A",
            "Category": props("Category").Value if props("Category") else "N/A",
            "Content Status": "N/A",
            "Identifier": "N/A",
            "Language": "N/A",
            "Revision": str(props("Revision Number").Value) if props("Revision Number") else "N/A"
        }

        prs.Close()
        app.Quit()
        pythoncom.CoUninitialize()  # 释放 COM 环境
        logger.info(f"Metadata extracted successfully from {file_path}")
        return metadata

    except Exception as e:
        logger.error(f"Error extracting metadata from legacy PPT format {file_path}: {e}")
        if 'app' in locals():
            app.Quit()
        pythoncom.CoUninitialize()  # 确保异常时也释放 COM
        return {"Error": str(e)}