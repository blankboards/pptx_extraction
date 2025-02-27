from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import logging
from datetime import datetime

# 配置日志
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

def extract_text_from_ppt(file_path):
    """
    从 PPT 文件中提取文本，保留层次结构并支持多种形状类型。

    Args:
        file_path (str): PPT 文件路径。

    Returns:
        list: 包含每张幻灯片文本的列表，格式为分隔符标记的字符串。
    """
    try:
        text_output = []
        presentation = Presentation(file_path)
        logger.info(f"Extracting text from PPT: {file_path}")

        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            logger.debug(f"Processing Slide {slide_number}")

            # 提取标题（如果有）
            if slide.shapes.title:
                slide_text.append(f"Title: {slide.shapes.title.text.strip()}")

            # 遍历所有形状
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # 处理文本框
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_text = "".join(run.text for run in paragraph.runs if run.text.strip())
                        if paragraph_text:
                            # 根据层级添加前缀（粗略判断）
                            level = paragraph.level if hasattr(paragraph, 'level') else 0
                            prefix = "  " * level + ("▪ " if level > 0 else "")
                            slide_text.append(f"{prefix}{paragraph_text}")
                
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    # 处理表格
                    table_text = []
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            table_text.append(" | ".join(row_text))
                    if table_text:
                        slide_text.append("Table:\n" + "\n".join(table_text))
                
                elif shape.has_chart:
                    # 处理图表（简单提取标题）
                    if shape.chart.chart_title.has_text_frame:
                        slide_text.append(f"Chart Title: {shape.chart.chart_title.text_frame.text.strip()}")

            # 处理注释（如果有）
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_text.append(f"Notes:\n{notes_text}")

            # 清理并格式化输出
            if slide_text:
                cleaned_text = "\n".join(line for line in slide_text if line.strip())
                text_output.append(f"@@@Slide_{slide_number}@@@\n{cleaned_text}")
            else:
                text_output.append(f"@@@Slide_{slide_number}@@@\n[No text content]")

        logger.info(f"Extracted text from {len(text_output)} slides.")
        return text_output

    except Exception as e:
        logger.error(f"Failed to extract text from {file_path}: {e}")
        return [f"Error: Unable to process {file_path}"]

def extract_metadata(file_path):
    """
    从 PPT 文件中提取元数据，包含更多核心属性。

    Args:
        file_path (str): PPT 文件路径。

    Returns:
        dict: 包含元数据的字典。
    """
    try:
        presentation = Presentation(file_path)
        props = presentation.core_properties
        logger.info(f"Extracting metadata from PPT: {file_path}")

        metadata = {
            "Title": props.title or "N/A",
            "Author": props.author or "N/A",
            "Subject": props.subject or "N/A",
            "Keywords": props.keywords or "N/A",
            "Comments": props.comments or "N/A",
            "Last Modified By": props.last_modified_by or "N/A",
            "Created": props.created.strftime('%Y-%m-%d %H:%M:%S') if props.created else "N/A",
            "Modified": props.modified.strftime('%Y-%m-%d %H:%M:%S') if props.modified else "N/A",
            "Category": props.category or "N/A",
            "Content Status": props.content_status or "N/A",
            "Identifier": props.identifier or "N/A",
            "Language": props.language or "N/A",
            "Revision": str(props.revision) if props.revision else "N/A"
        }

        logger.info("Metadata extraction completed.")
        return metadata

    except Exception as e:
        logger.error(f"Failed to extract metadata from {file_path}: {e}")
        return {"Error": f"Unable to process {file_path}"}

# 测试代码
if __name__ == "__main__":
    sample_ppt = "path/to/your/sample.pptx"
    text_result = extract_text_from_ppt(sample_ppt)
    metadata_result = extract_metadata(sample_ppt)
    
    print("Metadata:")
    for key, value in metadata_result.items():
        print(f"{key}: {value}")
    print("\nText:")
    for slide in text_result:
        print(slide)