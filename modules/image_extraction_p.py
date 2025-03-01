from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import logging
from paddleocr import PaddleOCR
from modules.config import OUTPUT_DIR
import win32com.client
import pythoncom

# 配置日志
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

def ensure_dir(directory):
    """确保目录存在，若不存在则创建"""
    if not os.path.exists(directory):
        os.makedirs(directory)
        logger.info(f"Created directory: {directory}")

def extract_images_from_ppt_paddleocr(file_path, output_dir=OUTPUT_DIR, output_format="text", use_gpu=False):
    """
    从 PPT 文件中提取图片，并使用 PaddleOCR 识别图片中的文本。

    Args:
        file_path (str): PPT 文件路径。
        output_dir (str): 输出目录路径，默认为 OUTPUT_DIR。
        output_format (str): 输出格式，可选 "text"（默认）或 "json"。
        use_gpu (bool): 是否使用 GPU 加速 OCR，默认 False。

    Returns:
        list: 包含每张图片识别文本的列表。
    """
    try:
        # 判断文件格式
        ext = os.path.splitext(file_path.lower())[1]
        if ext != '.pptx':
            return extract_images_from_ppt_legacy(file_path, output_dir, output_format, use_gpu)

        image_texts = []
        presentation = Presentation(file_path)
        logger.info(f"Processing PPT file: {file_path}")

        # 初始化 PaddleOCR
        ocr = PaddleOCR(use_angle_cls=True, lang='ch', use_gpu=use_gpu)

        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_folder = os.path.join(output_dir, f"slide_{slide_number}", "image")
            ensure_dir(slide_folder)
            image_index = 1

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    image_name = f"image_{image_index}.jpg"
                    image_path = os.path.join(slide_folder, image_name)

                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    logger.debug(f"Saved image: {image_path}")

                    if contains_text(image_path, ocr):
                        text = process_image_for_ocr(image_path, ocr)
                        if text:
                            text_entry = {
                                "slide": slide_number,
                                "image": image_index,
                                "text": text
                            }
                            if output_format == "json":
                                image_texts.append(text_entry)
                            else:
                                image_texts.append(f"\nSlide {slide_number}, Image {image_index} Text:\n{text}")

                            text_file_path = os.path.join(slide_folder, f"slide_{slide_number}_image_{image_index}_text.txt")
                            with open(text_file_path, "w", encoding="utf-8") as f:
                                f.write(text)
                            logger.info(f"Extracted text from {image_name}: {text[:50]}...")
                        else:
                            logger.debug(f"No readable text in {image_name}")
                    else:
                        logger.debug(f"No text detected in {image_name}")
                    image_index += 1

        logger.info(f"Completed processing {file_path}. Extracted {len(image_texts)} image texts.")
        return image_texts

    except Exception as e:
        logger.error(f"Failed to process PPT file {file_path}: {e}")
        return []

def process_image_for_ocr(image_path, ocr):
    """使用指定的 PaddleOCR 实例处理图片并提取文本"""
    try:
        result = ocr.ocr(image_path, cls=True)
        if not result or not result[0]:
            return ""

        texts = []
        for line in result[0]:
            if line and len(line[1][0]) >= 2 and line[1][1] > 0.5:
                texts.append(line[1][0])
        full_text = '\n'.join(texts).strip()
        return full_text if full_text else ""

    except Exception as e:
        logger.error(f"OCR processing failed for {image_path}: {e}")
        return ""

def contains_text(image_path, ocr):
    """使用指定的 PaddleOCR 实例判断图片是否含有文本"""
    try:
        result = ocr.ocr(image_path, cls=True)
        if result and result[0]:
            valid_texts = [
                line for line in result[0]
                if line and len(line[1][0]) >= 2 and line[1][1] > 0.5
            ]
            return len(valid_texts) > 0
        return False

    except Exception as e:
        logger.error(f"Text detection failed for {image_path}: {e}")
        return False

def extract_images_from_ppt_legacy(file_path, output_dir=OUTPUT_DIR, output_format="text", use_gpu=False):
    """
    从非 PPTX 格式的文件（如 .ppt, .pot, .pps）中提取图片并识别文本。

    Args:
        file_path (str): PPT 文件路径。
        output_dir (str): 输出目录路径，默认为 OUTPUT_DIR。
        output_format (str): 输出格式，可选 "text"（默认）或 "json"。
        use_gpu (bool): 是否使用 GPU 加速 OCR，默认 False。

    Returns:
        list: 包含每张图片识别文本的列表，与 extract_images_from_ppt_paddleocr 输出格式一致。
    """
    try:
        pythoncom.CoInitialize()
        image_texts = []
        app = win32com.client.Dispatch("PowerPoint.Application")
        prs = app.Presentations.Open(file_path, WithWindow=False)
        logger.info(f"Processing legacy PPT file: {file_path}")

        # 初始化 PaddleOCR
        ocr = PaddleOCR(use_angle_cls=True, lang='ch', use_gpu=use_gpu)

        for slide_number, slide in enumerate(prs.Slides, start=1):
            slide_folder = os.path.join(output_dir, f"slide_{slide_number}", "image")
            ensure_dir(slide_folder)
            image_index = 1

            for shape in slide.Shapes:
                if shape.Type == 13:
                    image_name = f"image_{image_index}.jpg"
                    image_path = os.path.join(slide_folder, image_name)
                    shape.Export(image_path, 2)
                    logger.debug(f"Saved image: {image_path}")

                    if contains_text(image_path, ocr):
                        text = process_image_for_ocr(image_path, ocr)
                        if text:
                            text_entry = {
                                "slide": slide_number,
                                "image": image_index,
                                "text": text
                            }
                            if output_format == "json":
                                image_texts.append(text_entry)
                            else:
                                image_texts.append(f"\nSlide {slide_number}, Image {image_index} Text:\n{text}")

                            text_file_path = os.path.join(slide_folder, f"slide_{slide_number}_image_{image_index}_text.txt")
                            with open(text_file_path, "w", encoding="utf-8") as f:
                                f.write(text)
                            logger.info(f"Extracted text from {image_name}: {text[:50]}...")
                        else:
                            logger.debug(f"No readable text in {image_name}")
                    else:
                        logger.debug(f"No text detected in {image_name}")
                    image_index += 1

            if image_index == 1:
                image_name = f"image_1.jpg"
                image_path = os.path.join(slide_folder, image_name)
                slide.Export(image_path, "JPG")
                logger.debug(f"Saved slide image: {image_path}")

                if contains_text(image_path, ocr):
                    text = process_image_for_ocr(image_path, ocr)
                    if text:
                        text_entry = {
                            "slide": slide_number,
                            "image": 1,
                            "text": text
                        }
                        if output_format == "json":
                            image_texts.append(text_entry)
                        else:
                            image_texts.append(f"\nSlide {slide_number}, Image 1 Text:\n{text}")

                        text_file_path = os.path.join(slide_folder, f"slide_{slide_number}_image_1_text.txt")
                        with open(text_file_path, "w", encoding="utf-8") as f:
                            f.write(text)
                        logger.info(f"Extracted text from {image_name}: {text[:50]}...")
                    else:
                        logger.debug(f"No readable text in slide {slide_number}")
                else:
                    logger.debug(f"No text detected in slide {slide_number}")

        prs.Close()
        app.Quit()
        pythoncom.CoUninitialize()
        logger.info(f"Completed processing {file_path}. Extracted {len(image_texts)} image texts.")
        return image_texts

    except Exception as e:
        logger.error(f"Failed to process legacy PPT file {file_path}: {e}")
        if 'app' in locals():
            app.Quit()
        pythoncom.CoUninitialize()
        return []